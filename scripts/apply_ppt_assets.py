#!/usr/bin/env python3
"""Inject visual assets and simulation screenshots into agents.pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT_PATH = ROOT / "agents.pptx"


def add_framed_image(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    """Add a subtle card frame + picture to a slide."""
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left - 0.06),
        Inches(top - 0.06),
        Inches(width + 0.12),
        Inches(height + 0.12),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(20, 27, 44)
    frame.line.color.rgb = RGBColor(70, 92, 132)
    frame.line.width = Pt(1.1)
    frame.shadow.inherit = False
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def main() -> int:
    prs = Presentation(str(PPT_PATH))

    assets = {
        "mcp_mesh": ROOT / "presentation_assets" / "custom" / "mcp_mesh.png",
        "planning_loop": ROOT / "presentation_assets" / "custom" / "planning_loop.png",
        "robotics_arm": ROOT / "presentation_assets" / "custom" / "robotics_arm.png",
        "sim1": ROOT / "robotics_maze" / "results" / "screenshots_mujoco" / "sim_mujoco_1_astar.png",
        "sim2": ROOT / "robotics_maze" / "results" / "screenshots_mujoco" / "sim_mujoco_2_weighted_astar.png",
        "sim3": ROOT / "robotics_maze" / "results" / "screenshots_mujoco" / "sim_mujoco_3_fringe_search.png",
    }

    for path in assets.values():
        if not path.exists():
            raise FileNotFoundError(f"Missing asset: {path}")

    placements = [
        # slide numbers are 1-based
        (6, "mcp_mesh", 7.0, 1.6, 2.7, 2.0),
        (18, "planning_loop", 6.8, 1.8, 2.8, 2.0),
        (32, "mcp_mesh", 6.9, 1.7, 2.8, 2.1),
        (35, "planning_loop", 6.9, 1.9, 2.8, 2.0),
        (37, "sim3", 5.4, 3.0, 4.4, 2.7),
        (38, "robotics_arm", 6.9, 1.6, 2.8, 2.1),
        (39, "sim1", 0.6, 3.0, 4.4, 2.7),
        (40, "sim2", 5.0, 3.0, 4.4, 2.7),
    ]
    for slide_number, asset_key, left, top, width, height in placements:
        slide_index = slide_number - 1
        if slide_index >= len(prs.slides):
            print(f"skip: slide {slide_number} missing (deck has {len(prs.slides)} slides)")
            continue
        add_framed_image(
            prs.slides[slide_index],
            assets[asset_key],
            left=left,
            top=top,
            width=width,
            height=height,
        )

    prs.save(str(PPT_PATH))
    print(f"updated: {PPT_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
