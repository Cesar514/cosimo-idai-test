#!/usr/bin/env python3
"""Apply full-deck polish: remove deprecated topics and ensure image + refs on each slide."""

from __future__ import annotations

import json
import tempfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Iterable
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT_PATH = ROOT / "agents.pptx"
IMAGE_MAP_PATH = ROOT / "presentation_assets" / "slide_image_map.json"
REF_MAP_PATH = ROOT / "presentation_assets" / "slide_references.json"
REMOVED_TOPIC_KEYWORDS = ("adoption", "closing")
SUPPORTED_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff"}
XML_NS_EP = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
XML_NS_VT = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"
NS = {"ep": XML_NS_EP, "vt": XML_NS_VT}


def iter_candidate_fallback_images() -> list[Path]:
    patterns = [
        ROOT / "presentation_assets" / "custom",
        ROOT / "robotics_maze" / "results" / "screenshots",
        ROOT / "robotics_maze" / "results" / "screenshots_mujoco",
        ROOT / "robotics_maze" / "testing" / "screenshots",
    ]
    files: list[Path] = []
    for base in patterns:
        if not base.exists():
            continue
        for suffix in ("*.png", "*.jpg", "*.jpeg"):
            files.extend(sorted(base.glob(suffix)))
    # Deduplicate while preserving order.
    seen: set[Path] = set()
    unique: list[Path] = []
    for path in files:
        resolved = path.resolve()
        if resolved in seen:
            continue
        seen.add(resolved)
        unique.append(resolved)
    return unique


def normalize_mapped_images(image_paths: Iterable[str]) -> list[Path]:
    normalized: list[Path] = []
    for raw in image_paths:
        path = Path(raw).expanduser()
        if not path.is_absolute():
            path = (ROOT / path).resolve()
        if path.suffix.lower() == ".svg":
            png_candidate = path.with_suffix(".png")
            if png_candidate.exists():
                path = png_candidate
            else:
                continue
        if path.suffix.lower() not in SUPPORTED_IMAGE_SUFFIXES:
            continue
        if path.exists():
            normalized.append(path)
    return normalized


def delete_slide(prs: Presentation, slide_index_zero_based: int) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001
    slide_ids = list(slide_id_list)
    slide_id = slide_ids[slide_index_zero_based]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide_id)


def strip_deprecated_paragraphs(text_frame) -> bool:
    removed_any = False
    to_remove = []
    for paragraph in text_frame.paragraphs:
        text = paragraph.text or ""
        lowered = text.lower()
        if any(keyword in lowered for keyword in REMOVED_TOPIC_KEYWORDS):
            to_remove.append(paragraph)
    for paragraph in to_remove:
        paragraph._p.getparent().remove(paragraph._p)  # noqa: SLF001
        removed_any = True
    if removed_any and not text_frame.paragraphs:
        text_frame.clear()
    return removed_any


def remove_text_mentions(slide) -> int:
    removed_count = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if strip_deprecated_paragraphs(shape.text_frame):
            removed_count += 1
    return removed_count


def align_agenda_timing(slide) -> bool:
    """Keep Agenda duration coherent after removing deprecated wrap-up bullets."""
    if not slide.shapes.title or not slide.shapes.title.text:
        return False
    title_text = slide.shapes.title.text.strip()
    if "agenda" not in title_text.lower() or "60" not in title_text:
        return False

    changed = False
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = (paragraph.text or "").strip()
            if text.startswith("11:45-11:55:"):
                paragraph.text = text.replace("11:45-11:55:", "11:45-12:00:", 1)
                changed = True
    return changed


def slide_title(slide, index_one_based: int) -> str:
    if slide.shapes.title and slide.shapes.title.text:
        cleaned = slide.shapes.title.text.strip()
        if cleaned:
            return cleaned
    return f"Slide {index_one_based}"


def sync_extended_properties(ppt_path: Path, slide_titles: list[str]) -> bool:
    with ZipFile(ppt_path) as src:
        entries = {name: src.read(name) for name in src.namelist()}

    app_key = "docProps/app.xml"
    if app_key not in entries:
        return False

    root = ET.fromstring(entries[app_key])
    changed = False

    slides_elem = root.find("ep:Slides", NS)
    if slides_elem is not None and slides_elem.text != str(len(slide_titles)):
        slides_elem.text = str(len(slide_titles))
        changed = True

    heading_variants = root.findall("ep:HeadingPairs/vt:vector/vt:variant", NS)
    non_slide_title_count = 0
    for idx in range(0, len(heading_variants), 2):
        if idx + 1 >= len(heading_variants):
            break
        name_elem = heading_variants[idx].find("vt:lpstr", NS)
        count_elem = heading_variants[idx + 1].find("vt:i4", NS)
        if name_elem is None or count_elem is None or count_elem.text is None:
            continue
        name = (name_elem.text or "").strip()
        if name == "Slide Titles":
            if count_elem.text != str(len(slide_titles)):
                count_elem.text = str(len(slide_titles))
                changed = True
        else:
            non_slide_title_count += int(count_elem.text)

    titles_vector = root.find("ep:TitlesOfParts/vt:vector", NS)
    if titles_vector is not None:
        existing_titles = [elem.text or "" for elem in titles_vector.findall("vt:lpstr", NS)]
        prefix = existing_titles[:non_slide_title_count]
        new_titles = prefix + slide_titles
        if existing_titles != new_titles:
            for child in list(titles_vector):
                titles_vector.remove(child)
            for text in new_titles:
                lpstr = ET.SubElement(titles_vector, f"{{{XML_NS_VT}}}lpstr")
                lpstr.text = text
            titles_vector.set("size", str(len(new_titles)))
            changed = True

    if not changed:
        return False

    entries[app_key] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    with tempfile.NamedTemporaryFile("wb", suffix=".pptx", dir=ppt_path.parent, delete=False) as tmp:
        tmp_path = Path(tmp.name)
    try:
        with ZipFile(tmp_path, "w", compression=ZIP_DEFLATED) as out:
            for name, payload in entries.items():
                out.writestr(name, payload)
        tmp_path.replace(ppt_path)
    finally:
        if tmp_path.exists():
            tmp_path.unlink(missing_ok=True)
    return True


def slide_has_picture(slide) -> bool:
    return count_integrated_images(slide) > 0


def count_integrated_images(slide) -> int:
    return sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)


def build_image_validation_rows(prs: Presentation) -> list[tuple[int, int]]:
    rows: list[tuple[int, int]] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        rows.append((slide_number, count_integrated_images(slide)))
    return rows


def add_framed_image(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left - 0.05),
        Inches(top - 0.05),
        Inches(width + 0.10),
        Inches(height + 0.10),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(20, 27, 44)
    frame.line.color.rgb = RGBColor(72, 95, 137)
    frame.line.width = Pt(1.0)
    frame.shadow.inherit = False
    slide.shapes.add_picture(
        str(image_path),
        Inches(left),
        Inches(top),
        width=Inches(width),
        height=Inches(height),
    )


def has_reference_footer(slide, slide_height_emu: int) -> bool:
    footer_threshold = slide_height_emu - Inches(0.9)
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if shape.top < footer_threshold:
            continue
        text = shape.text_frame.text or ""
        lowered = text.lower()
        if "http" in lowered or "refs:" in lowered:
            return True
    return False


def add_reference_footer(slide, refs: list[str]) -> None:
    footer = slide.shapes.add_textbox(Inches(0.28), Inches(7.06), Inches(12.78), Inches(0.34))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Refs: " + " | ".join(refs[:3])
    p.font.size = Pt(6.5)
    p.font.color.rgb = RGBColor(117, 123, 136)


def main() -> int:
    prs = Presentation(str(PPT_PATH))
    image_map = json.loads(IMAGE_MAP_PATH.read_text())
    refs_map = json.loads(REF_MAP_PATH.read_text())
    fallback_images = iter_candidate_fallback_images()
    if not fallback_images:
        raise RuntimeError("No fallback images found.")

    # Remove any whole slides still titled with removed topics.
    slide_indices_to_delete = []
    removed_slide_titles: list[str] = []
    for idx, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip().lower()
        if title and any(keyword in title for keyword in REMOVED_TOPIC_KEYWORDS):
            slide_indices_to_delete.append(idx)
            removed_slide_titles.append(title)
    for idx in reversed(slide_indices_to_delete):
        delete_slide(prs, idx)

    fallback_idx = 0
    added_images = 0
    added_refs = 0
    scrubbed_shape_mentions = 0
    agenda_alignment_count = 0

    for slide_number, slide in enumerate(prs.slides, start=1):
        scrubbed_shape_mentions += remove_text_mentions(slide)
        if align_agenda_timing(slide):
            agenda_alignment_count += 1

        if not slide_has_picture(slide):
            mapped = normalize_mapped_images(image_map.get(str(slide_number), []))
            if mapped:
                image = mapped[0]
            else:
                image = fallback_images[fallback_idx % len(fallback_images)]
                fallback_idx += 1
            add_framed_image(slide, image, left=9.06, top=4.88, width=3.92, height=2.32)
            added_images += 1

        if not has_reference_footer(slide, prs.slide_height):
            refs = refs_map.get(str(slide_number), [])
            if not refs:
                refs = refs_map.get("1", [])
            add_reference_footer(slide, refs)
            added_refs += 1

    validation_rows = build_image_validation_rows(prs)
    missing_images = [slide_number for slide_number, image_count in validation_rows if image_count < 1]
    if missing_images:
        raise RuntimeError(f"Slides without integrated images: {missing_images}")

    final_slide_titles = [slide_title(slide, idx) for idx, slide in enumerate(prs.slides, start=1)]
    prs.save(str(PPT_PATH))
    synced_properties = sync_extended_properties(PPT_PATH, final_slide_titles)

    print(f"updated: {PPT_PATH}")
    print(f"slides_after={len(prs.slides)} added_images={added_images} added_refs={added_refs}")
    print(f"removed_titled_slides={len(removed_slide_titles)}")
    if removed_slide_titles:
        print("removed_slide_titles=" + " | ".join(removed_slide_titles))
    print(f"slides_with_text_scrubbed={scrubbed_shape_mentions}")
    print(f"agenda_timing_aligned={agenda_alignment_count}")
    print(f"extended_properties_synced={int(synced_properties)}")
    print("image_validation_report_start")
    for slide_number, image_count in validation_rows:
        print(f"slide={slide_number:02d} image_count={image_count}")
    print("image_validation_report_end")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
